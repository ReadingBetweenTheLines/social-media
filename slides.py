import os
import sys
from typing import List, Optional
from pydantic import BaseModel, Field
from openai import OpenAI
from docx import Document
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

# =====================================================================
# 1. STRICT DATA SCHEMAS (Pydantic Validation)
# =====================================================================
class SlideItem(BaseModel):
    title: str = Field(..., description="Short, high-impact slide header (max 6 words).")
    bullets: List[str] = Field(..., description="3 to 5 bite-sized bullet points. No long paragraphs.")
    speaker_notes: Optional[str] = Field(None, description="Contextual guidance or script for the presenter.")

class PresentationStructure(BaseModel):
    title: str = Field(..., description="The main, clear title for the presentation slide deck.")
    subtitle: str = Field(..., description="A short subtitle providing context or summary.")
    slides: List[SlideItem] = Field(..., description="Chronological listing of text content slides.")


# =====================================================================
# 2. CORE CONVERTER ENGINE
# =====================================================================
class RobustDocConverter:
    def __init__(self):
        # Initializes client with system-level environment token automatically
        if not os.environ.get("OPENAI_API_KEY"):
            raise EnvironmentError("Missing 'OPENAI_API_KEY' environment variable.")
        self.client = OpenAI()

    def extract_document_text(self, filepath: str) -> str:
        """Parses document variations safely and returns continuous text strings."""
        if not os.path.exists(filepath):
            raise FileNotFoundError(f"Target file path not found: {filepath}")

        ext = os.path.splitext(filepath)[-1].lower()
        
        try:
            if ext in [".txt", ".md"]:
                with open(filepath, "r", encoding="utf-8") as f:
                    return f.read().strip()
            elif ext == ".docx":
                doc = Document(filepath)
                return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
            else:
                raise ValueError(f"Unsupported file format extension: {ext}")
        except Exception as e:
            raise RuntimeError(f"Failed parsing file contents securely: {str(e)}")

    def analyze_and_structure_content(self, raw_text: str) -> PresentationStructure:
        """Leverages Structured Outputs to generate predictive schema blocks without formatting drops."""
        print("[*] Processing raw text via AI Structural Parsing Engine...")
        
        system_instruction = (
            "You are a professional presentation design backend engine. Your task is to extract "
            "the core themes from data payloads and format them into sequence slides. "
            "Strip out fluff, condense paragraphs into highly technical or action-oriented bullets, "
            "and map a cohesive flow. You must rigorously adhere to the JSON output formatting rules."
        )

        # Utilizing the strict structural parsing wrapper from modern SDK releases
        completion = self.client.beta.chat.completions.parse(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": system_instruction},
                {"role": "user", "content": f"Transform this source content into a slide layout structure:\n\n{raw_text}"}
            ],
            response_format=PresentationStructure,
        )
        return completion.choices[0].message.parsed

    def build_presentation_file(self, structure: PresentationStructure, output_path: str, template_path: str = None):
        """Builds a .pptx file applying structural formatting bounds to avoid text overflows."""
        print(f"[*] Building slide layout layers to: {output_path}")
        
        # Load custom corporate layout master if provided, else spawn base templates
        prs = Presentation(template_path) if template_path else Presentation()
        
        # Layout index configurations maps (Standard default sets)
        TITLE_LAYOUT = prs.slide_layouts[0]
        CONTENT_LAYOUT = prs.slide_layouts[1]

        # --- SLIDE 1: Title Layout Placement ---
        title_slide = prs.slides.add_slide(TITLE_LAYOUT)
        title_slide.shapes.title.text = structure.title
        if hasattr(title_slide.placeholders, 'length') and len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = structure.subtitle or ""

        # --- SLIDE 2+: Core Content Iteration ---
        for item in structure.slides:
            slide = prs.slides.add_slide(CONTENT_LAYOUT)
            
            # Configure Header Properties Safely
            title_shape = slide.shapes.title
            title_shape.text = item.title
            title_shape.text_frame.word_wrap = True
            
            # Extract Content Frame Placements
            content_placeholder = slide.placeholders[1]
            tf = content_placeholder.text_frame
            tf.word_wrap = True  # Vital safeguard parameter ensuring automatic wrapper containment
            
            # Inject bullet point matrices defensively adjusting font dynamics 
            # to accommodate varying sentence blocks without overflow.
            for idx, bullet_text in enumerate(item.bullets):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.text = bullet_text
                p.level = 0
                
                # Dynamic visual spacing tuning thresholds
                if len(item.bullets) >= 5 or len(bullet_text) > 120:
                    p.font.size = Pt(14)
                else:
                    p.font.size = Pt(18)
                    
            # Populate presentation voice tracks/speaker guidance if present
            if item.speaker_notes:
                slide.notes_slide.notes_text_frame.text = item.speaker_notes

        # Commit presentation layers to physical secondary storage arrays
        prs.save(output_path)
        print(f"[✓] System workflow executed cleanly. Output generated inside: {output_path}")


# =====================================================================
# 3. OPERATION ENTRY ENGINE CLI
# =====================================================================
if __name__ == "__main__":
    # Test script parameter validation requirements
    if len(sys.argv) < 2:
        print("Usage Error: python doc_to_slides.py <path_to_source_document> [optional_output_path.pptx]")
        sys.exit(1)

    input_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else "generated_deck.pptx"

    try:
        converter = RobustDocConverter()
        
        # Step 1: Safely ingest document structure maps
        extracted_text = converter.extract_document_text(input_file)
        print(f"[✓] Successfully ingested target asset file context ({len(extracted_text)} characters).")
        
        # Step 2: Extract data payloads strictly inside structure schemas
        structured_data = converter.analyze_and_structure_content(extracted_text)
        
        # Step 3: Draw out layout renderings safely
        converter.build_presentation_file(structured_data, output_file)
        
    except Exception as error_context:
        print(f"[!] System Exception Caught During Production Pipeline Execution:\n    {str(error_context)}")
        sys.exit(1)